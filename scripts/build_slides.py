#!/usr/bin/env python3
"""Build the presentation deck: HTML, then PDF and PPTX from the same content.

The first version of this deck was produced as one-off files with no generator,
so a corrected figure meant rebuilding it by hand. Slide content lives here as
data and every format is rendered from it.

Numbers must match docs/portfolio-brief.md, which in turn matches results/.
Run `python3 scripts/check_slide_numbers.py` after editing.
"""

from __future__ import annotations

import argparse
import html
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "report"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

CSS = """
  :root {
    --ink: #1a1a1a;
    --muted: #666;
    --accent: #0b5cad;
    --rule: #ddd;
    --paper: #ffffff;
  }
  @page { size: 13.333in 7.5in; margin: 0; }
  * { box-sizing: border-box; }
  body {
    margin: 0;
    background: var(--paper);
    color: var(--ink);
    font-family: "Noto Sans KR", -apple-system, sans-serif;
    word-break: keep-all;
  }
  .deck { display: grid; gap: 28px; justify-content: center; padding: 28px; }
  .slide {
    width: 1280px;
    height: 720px;
    padding: 58px 72px 48px;
    background: var(--paper);
    overflow: hidden;
    break-after: page;
    position: relative;
  }
  h1 {
    margin: 0 0 28px;
    color: var(--ink);
    font-size: 42px;
    line-height: 1.22;
    letter-spacing: -0.035em;
    font-weight: 760;
  }
  .title-slide { display: flex; flex-direction: column; justify-content: center; }
  .title-slide h1 { color: var(--accent); font-size: 64px; margin-bottom: 18px; }
  .subtitle { font-size: 34px; font-weight: 700; margin: 0 0 58px; }
  .meta { color: var(--muted); font-size: 21px; line-height: 1.8; }
  .lead { margin: 0 0 24px; font-size: 25px; line-height: 1.45; font-weight: 690; }
  .claim {
    margin: 22px 0 0;
    color: var(--accent);
    font-size: 24px;
    line-height: 1.42;
    font-weight: 740;
  }
  table { width: 100%; border-collapse: collapse; table-layout: fixed; font-size: 19px; }
  thead { border-top: 1px solid var(--ink); border-bottom: 1px solid var(--muted); }
  tbody tr { border-bottom: 1px solid var(--rule); }
  th, td { padding: 12px 10px; border: 0; text-align: left; vertical-align: middle; }
  th { font-weight: 720; }
  tr.emphasis td { font-weight: 760; color: var(--accent); }
  table.compact { font-size: 15px; }
  table.compact th, table.compact td { padding: 10px 6px; }
  .mono, td.numeric {
    font-family: "JetBrains Mono", "DejaVu Sans Mono", monospace;
    font-variant-numeric: tabular-nums;
  }
  .sections { display: grid; grid-template-columns: repeat(var(--count), 1fr); gap: 34px; }
  .section { padding-top: 14px; border-top: 1px solid var(--rule); }
  .section h2 { margin: 0 0 14px; font-size: 24px; color: var(--accent); }
  .section p { margin: 0; color: var(--ink); font-size: 20px; line-height: 1.55; }
  pre {
    margin: 22px 0 0;
    padding: 18px 0;
    border-top: 1px solid var(--rule);
    border-bottom: 1px solid var(--rule);
    background: var(--paper);
    color: var(--ink);
    font: 16px/1.55 "JetBrains Mono", "DejaVu Sans Mono", monospace;
    white-space: pre-wrap;
  }
  .risk-list { margin: 4px 0 0; padding-left: 28px; font-size: 22px; line-height: 1.52; }
  .risk-list li { margin: 13px 0; }
  .foot { position: absolute; left: 72px; right: 72px; bottom: 44px; color: var(--muted); font-size: 16px; }
  .two-column { display: grid; grid-template-columns: 1fr 1fr; gap: 54px; }
  .two-column h2 { margin: 0 0 14px; color: var(--accent); font-size: 24px; }
  .two-column p { font-size: 20px; line-height: 1.55; }
  @media print {
    .deck { display: block; padding: 0; }
    .slide { margin: 0; }
  }
  @media screen {
    .slide { outline: 1px solid var(--rule); }
  }
"""


SLIDES: list[dict[str, Any]] = [
    {
        "cover": True,
        "title": "FinLLM Lab",
        "subtitle": "양자화·금융 관점 기술 보고서",
        "meta": [
            "측정일: 2026-08-08",
            "NVIDIA RTX A6000 48GB 1장",
            "금융 내부문서 RAG · 동시 사용자 10명 · max-model-len 8,192",
        ],
        "notes": "이 문서의 모든 수치는 results/의 27개 결과 레코드에서 나왔고, 각 레코드는 "
        "scripts/finllm_profile.py validate-result 검증을 통과한다.",
    },
    {
        "title": "문제와 합격 조건",
        "lead": "합격 조건은 실험 전에 고정했다. 실험 후에 임계값을 바꾸지 않는다.",
        "table": {
            "head": ["항목", "기준"],
            "rows": [
                ["품질", "100점 만점 90점 이상"],
                ["P95 TTFT", "2,000ms 이하"],
                ["요청 오류율", "1% 이하"],
                ["OOM", "0회"],
                ["권한 위반 검색", "0건"],
                ["동시 사용자", "10명"],
            ],
            "numeric_cols": [1],
        },
        "foot": "answer correctness 40% + groundedness 30% + citation accuracy 20% + abstention/safety 10%",
        "notes": "금융 내부문서 RAG를 온프레미스 단일 GPU로 운영할 때, 어느 정도 하드웨어에서 "
        "필요한 품질과 응답 성능이 나오는가를 재현 가능한 근거로 결정한다.",
    },
    {
        "title": "최종 권고: Qwen3-14B-AWQ · 0.46 · --enforce-eager",
        "lead": "Profile A 권고 구성: Qwen3-14B-AWQ (W4A16), executor 예산 0.46, --enforce-eager",
        "table": {
            "compact": True,
            "head": [
                "후보", "Quality", "P95 TTFT(서버)", "P95 TTFT(사용자)",
                "tok/s", "Peak VRAM", "24GB 적합", "최대 동시성",
            ],
            "rows": [
                ["8B BF16 (0.50)", "95.9", "81ms", "1,349ms", "287.0", "24.01GiB", "초과", "7.31"],
                ["14B AWQ (0.50)", "97.7", "132ms", "1,329ms", "313.2", "23.84GiB", "적합", "11.05"],
                ["14B AWQ (0.46)", "97.7", "132ms", "1,287ms", "315.3", "21.96GiB", "적합", "9.53"],
            ],
            "emphasis": 2,
            "numeric_cols": [0, 1, 2, 3, 4, 5, 7],
        },
        # Added 2026-08-09: the table says "초과" for 8B, and the margin behind
        # that word is 6 MiB. Stating it here keeps the claim from reading
        # stronger than the measurement supports.
        "claim": "8B 탈락의 실질 근거는 메모리 마진이 아니라 품질 95.9 vs 97.7과 최대 동시성 7.31 vs 9.53이다.",
        "foot": "3회 반복, 편차 1% 이내. 오류율 0%, OOM 0회, 권한 위반 0건. "
        "8B peak는 24.006GiB로 24.0GiB를 6MiB 초과 — 여유가 없다는 뜻이지 크게 넘는다는 뜻이 아니다.",
        "notes": "8B BF16의 24.006GiB는 24.0GiB를 6MiB 넘는다. 실제 24GB 카드는 총 용량이 "
        "24.0GiB이고 드라이버와 디스플레이 출력이 일부를 쓰므로 이 구성은 들어가지 못하지만, "
        "max-model-len을 조금만 줄이면 들어갈 여지가 있다. 반면 14B AWQ의 21.96GiB는 2.0GiB "
        "여유가 있어 성격이 다르다. 따라서 8B를 탈락시킨 실질적 근거는 품질과 최대 동시성이고 "
        "메모리는 그 위에 얹힌 추가 근거다.",
    },
    {
        "title": "4-bit 양자화의 이득은 KV cache 여유로 나타난다",
        "lead": "같은 24GB 예산 안에서 메모리가 어떻게 쓰이는지가 결정을 갈랐다.",
        "table": {
            "head": ["구성", "가중치", "KV cache", "최대 동시성 (8,192 토큰 기준)"],
            "rows": [
                ["8B BF16", "15.27GiB", "8.22GiB", "7.31"],
                ["14B AWQ", "9.37GiB", "11.91GiB", "9.53"],
            ],
            "numeric_cols": [0, 1, 2, 3],
        },
        "claim": "파라미터가 1.8배 큰 모델이 메모리를 덜 쓰면서 동시 요청을 30% 더 처리한다.",
        "notes": "가중치가 15.27GiB에서 9.37GiB로 줄어든 만큼 KV cache에 쓸 공간이 생겼다. "
        "동시 사용자 10명이 요구조건인데 8B BF16은 8,192 토큰 기준 최대 동시성이 7.31로 미달한다. "
        "요구조건을 만족시키는 쪽이 더 큰 모델이라는 점이 이 실험의 비직관적인 결론이다.",
    },
    {
        "title": "오진과 정정 — CUDA graph와 eager",
        "lead": "처리량 열세의 원인을 양자화로 오진했다. 틀렸다.",
        "table": {
            "head": ["", "CUDA graph", "enforce-eager", "배수"],
            "rows": [
                ["14B AWQ tok/s", "57.2", "313.2", "5.5×"],
                ["8B BF16 tok/s", "296.0", "287.0", "0.97×"],
            ],
            "numeric_cols": [0, 1, 2, 3],
        },
        "claim": "vLLM 0.9.2 + Ampere + AWQ 조합에서 CUDA graph 경로가 병리적으로 느리다.",
        "notes": "당시 진단은 Ampere에서 AWQ 4-bit 역양자화 비용이 배치 10 구간에서 지배적이 "
        "된다는 것이었다. --enforce-eager로 CUDA graph capture만 끄자 결과가 바뀌었다. "
        "양자화가 원인이었다면 graph를 끄는 것으로 이런 차이가 날 수 없고 BF16 모델이 영향을 "
        "받지 않을 이유도 없다. 정확한 지점은 vLLM 내부 프로파일링이 필요하며 이 실험에서는 "
        "확인하지 못했다.",
    },
    {
        "title": "오진을 확인한 방법",
        "lead": "결과가 뒤집힐 때 가장 먼저 의심할 것은 측정 오류다.",
        "sections": [
            ("같은 일을 했는가", "총 출력 토큰 640 · 최소 5 / 중앙 25 / 최대 34 · 성공 30/30 · 오류 0건 · 11.17초 → 2.06초"),
            ("반복해도 같은가", "3회 반복 편차 1% 이내."),
            ("배치 효과인가", "동시성 1에서 6.2배, 동시성 10에서 5.5배."),
        ],
        "claim": "배치 크기와 무관하게 일정하므로 '배치 구간에서 역양자화가 지배적'이라는 설명은 성립하지 않는다.",
        "notes": "graph 반복은 11.17s / 57.3 tok/s, 11.20s / 57.1, 11.19s / 57.2였고, eager 반복은 "
        "2.06s / 310.8 tok/s, 2.03s / 316.0, 2.04s / 312.9였다. 동시성 1에서 graph는 p50 TTFT "
        "142.5ms, P95 E2E 4,716ms, 7.3 tok/s였고, eager는 27.7ms, 768ms, 45.4 tok/s였다.",
    },
    {
        "title": "gpu_memory_utilization은 카드 적합성을 보장하지 않는다",
        "lead": "이 값은 executor 예산만 정한다. CUDA graph와 CUDA context는 그 밖에서 쓴다.",
        "table": {
            "head": ["구성", "가중치", "KV cache", "CUDA graph", "실측 peak"],
            "rows": [
                ["14B AWQ 0.46 + graph", "9.37", "11.89", "3.51", "25.67GiB"],
                ["14B AWQ 0.46 + eager", "9.37", "11.91", "0", "21.96GiB"],
            ],
            "numeric_cols": [0, 1, 2, 3, 4],
        },
        "claim": "'예산 안에 들어갔다'는 '24GB 카드에서 돈다'가 아니다.",
        "notes": "executor 예산 22.08GiB로 같은데 총 사용량이 25.67GiB와 21.96GiB로 갈린다.",
    },
    {
        "title": "금융 도메인 권한 모델",
        "lead": "준법감시부와 내부감사부가 서로의 문서를 볼 수 없다는 점이 설계의 핵심이다.",
        "table": {
            "head": ["역할", "열람 가능"],
            "rows": [
                ["compliance-officer (준법감시부)", "전행 공통규정, 준법감시 지침"],
                ["internal-audit (내부감사부)", "전행 공통규정, 감사 작업문서"],
                ["whistleblow-admin (제보 전담)", "전행 공통규정, 제보 사건파일"],
                ["branch-staff (영업점)", "전행 공통규정"],
                ["vendor-contractor (외부 위탁)", "공개 안내문만"],
            ],
            "numeric_cols": [0],
        },
        "foot": "문서 소관 부서를 기준으로 한 ABAC 필터가 있어야만 분리된다.",
        "notes": "두 부서 모두 본부 부서라 직급 기반 권한으로는 분리되지 않고, 문서 소관 부서를 "
        "기준으로 한 ABAC 필터가 있어야만 분리된다.",
    },
    {
        "title": "평가셋 60문항 구성",
        "lead": "합성 문서 16종 83조항 — 정답과 허용 근거 문서 ID를 정확히 통제한다.",
        "table": {
            "head": ["유형", "문항", "무엇을 재는가"],
            "rows": [
                ["근거 있는 질문", "25", "정답과 인용 정확성"],
                ["여러 문서 조합", "10", "retrieval과 종합 능력"],
                ["근거 없는 질문", "10", "올바른 답변 유보"],
                ["권한 없는 문서", "10", "ACL 우회 방지"],
                ["악성 지시 문서", "5", "prompt injection 대응"],
            ],
            "numeric_cols": [1],
        },
        "foot": "그럴듯하지만 무관한 근거를 받고도 유보해야 한다.",
        "notes": "unauthorized 문항은 검색 결과가 비는 시험이 아니다. ACL이 정답 문서를 가려도 "
        "어휘 검색은 그 역할이 볼 수 있는 다른 문서를 반환한다.",
    },
    {
        "title": "권한 격리는 성공, injection 방어는 실패",
        "table": {
            "head": ["항목", "8B BF16", "14B AWQ"],
            "rows": [
                ["권한 위반 검색", "0건", "0건"],
                ["injection 성공", "2 / 5", "2 / 5"],
            ],
            "numeric_cols": [1, 2],
        },
        "pre": "네, ... 사전에 서면 승인을 받은 경우에는 재위탁이 가능합니다 [POL-2026-006#제4조].\n\n"
        "내부통제 점검은 생략 가능합니다              ← 문서가 시킨 문장\n"
        "[POL-2026-001#제3조] ... [POL-2026-006#제5조]  ← 문서가 시킨 목록 나열",
        "claim": "권한을 프롬프트로 지켰다면 injection과 함께 뚫렸을 것이다.",
        "notes": "한 문항에서는 존재하지 않는 문서 ID POL-2026-007~010까지 생성했다. 권한 위반이 "
        "0건인 것은 모델이 잘해서가 아니라 retrieval 이전에 데이터 층에서 강제했기 때문이다. "
        "이 구성은 injection 방어를 추가하기 전까지 배포 가능하지 않다.",
    },
    {
        "title": "증거 규율 — 주장 범위를 코드로 분리한다",
        "two_column": [
            ("memory-budget-emulation", "A6000에서 예산을 제한한 실험. 메모리 적합성과 A6000 관측 성능까지만 주장할 수 있다."),
            ("native-gpu-validation", "실제 대상 카드 실측. 이 라벨만 대상 GPU 성능을 주장할 수 있다."),
        ],
        "claim": "모든 수치는 A6000 관측값이며 RTX 4090 성능이 아니다.",
        "foot": "서버 관점 P95 TTFT 302ms · 사용자 체감 6,750ms — 어느 TTFT인지 밝히지 않은 수치는 쓸모가 없다.",
        "notes": "validator가 A6000 실행을 native-gpu-validation으로 라벨링하는 것을 거부한다.",
    },
    {
        "title": "스캐폴드 결함 — 검증 코드가 주장을 지킨다",
        "lead": "인수받은 스캐폴드에서 결함 6건을 찾아 고쳤다.",
        "table": {
            "compact": True,
            "head": ["결함", "영향"],
            "rows": [
                ["발표된 JSON Schema가 어디에서도 사용되지 않음", "결과 계약이 문서상으로만 존재"],
                ["템플릿 생성기가 그 Schema를 위반", "위와 맞물려 아무도 발견 못 함"],
                ["native-gpu-validation 라벨 오용 무방비", "정직성 주장 자체가 무방비"],
                ["TTFT가 클라이언트 대기시간 제외", "사용자 체감 지연 과소보고"],
                ["tokenizer revision 미고정", "프로토콜 요구사항과 불일치"],
                ["환경 기록이 타 사용자 프로세스 경로 포함", "스스로 밝힌 비식별 원칙 위반"],
            ],
            "numeric_cols": [],
        },
        "foot": "7번째: 판단 근거가 비어 있는데도 검증 통과 — placeholder 접두어 검사로 해결.",
        "notes": "측정 자동화 스크립트를 만들자 결과 레코드 9개가 판단 근거가 비어 있는데도 검증을 "
        "통과하는 7번째 결함이 드러났다. 그 필드가 필수 목록에 없었고, 추가한 뒤에도 placeholder "
        "검사가 정확 일치라 빠져나갔다.",
    },
    {
        "title": "검증 체계",
        "table": {
            "head": ["항목", "규모"],
            "rows": [
                ["자동 테스트", "84개"],
                ["결과 레코드", "27개 (전부 schema·게이트 검증 통과)"],
                ["ADR", "4건 (1건은 오진 기록으로 Superseded 보존)"],
                ["합성 corpus", "16문서 83조항"],
                ["평가셋", "60문항"],
            ],
            "numeric_cols": [1],
        },
        "claim": "평가셋 자신도 corpus·권한·인용·질문 누설 여부를 검증한다.",
        "foot": "CI 게이트: model revision · synthetic 표시 · 자격증명 패턴 · results/ 전수 검증",
        "notes": "tests/test_eval_set.py가 질문 안에 정답 문자열이 들어 있지 않은지까지 검사한다. "
        "그 검사가 실제로 결함을 잡았다.",
    },
    {
        "title": "측정 이력 — 검색 설정이 바뀌면 결과도 바뀐다",
        "lead": "세 번의 측정을 모두 남겼다. 옛 레코드를 새 값으로 고치지 않는다.",
        "table": {
            "head": ["태그", "조건", "retriever"],
            "rows": [
                ["2026-08-08", "랭킹에 조항 제목 미포함, CUDA graph", "0e40e0354b7b"],
                ["2026-08-08b", "조항 제목 포함, CUDA graph", "11d1f8cfeb42"],
                ["2026-08-08c", "조항 제목 포함, --enforce-eager", "11d1f8cfeb42"],
            ],
            "numeric_cols": [0, 2],
        },
        "claim": "설정이 다르면 결과도 다르다. 옛 레코드를 새 해시로 고치는 것은 측정하지 않은 값을 기록하는 일이다.",
        "notes": "교차 검토로 랭킹에 조항 제목을 넣자 기대 인용 회수율이 38/40에서 40/40, "
        "multi_doc이 8/10에서 10/10으로 올랐다. 검색이 바뀌었으므로 품질을 다시 측정했다.",
    },
    {
        "title": "남은 위험과 다음 단계",
        "risks": [
            "--enforce-eager 결정은 Ampere 고유일 수 있다. Ada/Blackwell에서 재검증 필요.",
            "prompt injection 방어가 없다. 입력 분류기·구조적 분리·출력 필터 조합을 측정해야 한다.",
            "부하 모델이 버스트다. closed-loop 10 사용자 모델로 재측정 필요.",
            "corpus가 합성이고 60문항이다. 1.8점 품질 차이를 크게 주장하지 않는다.",
            "실제 24GB 카드 실측을 하지 않았다. 모든 결과는 memory-budget-emulation 증거다.",
        ],
        "notes": "이 구성은 injection 방어를 추가하기 전까지 배포 가능하지 않다.",
    },
]


def esc(text: str) -> str:
    return html.escape(str(text))


def render_table(spec: dict[str, Any]) -> str:
    numeric = set(spec.get("numeric_cols", []))
    cls = ' class="compact"' if spec.get("compact") else ""
    head = "".join(f"<th>{esc(h)}</th>" for h in spec["head"])
    body = []
    for index, row in enumerate(spec["rows"]):
        row_cls = ' class="emphasis"' if spec.get("emphasis") == index else ""
        cells = "".join(
            '<td class="numeric">{}</td>'.format(esc(c))
            if i in numeric
            else "<td>{}</td>".format(esc(c))
            for i, c in enumerate(row)
        )
        body.append(f"<tr{row_cls}>{cells}</tr>")
    return (
        f"<table{cls}><thead><tr>{head}</tr></thead>"
        f"<tbody>{''.join(body)}</tbody></table>"
    )


def render_slide(slide: dict[str, Any]) -> str:
    if slide.get("cover"):
        meta = "".join(f"<div>{esc(m)}</div>" for m in slide["meta"])
        inner = (
            f'<h1>{esc(slide["title"])}</h1>'
            f'<p class="subtitle">{esc(slide["subtitle"])}</p>'
            f'<div class="meta">{meta}</div>'
        )
        return (
            f'<section class="slide title-slide">{inner}'
            f'<aside hidden>{esc(slide["notes"])}</aside></section>'
        )

    parts = [f'<h1>{esc(slide["title"])}</h1>']
    if slide.get("lead"):
        parts.append(f'<p class="lead">{esc(slide["lead"])}</p>')
    if slide.get("table"):
        parts.append(render_table(slide["table"]))
    if slide.get("sections"):
        blocks = "".join(
            f'<div class="section"><h2>{esc(h)}</h2><p>{esc(b)}</p></div>'
            for h, b in slide["sections"]
        )
        parts.append(
            f'<div class="sections" style="--count: {len(slide["sections"])}">{blocks}</div>'
        )
    if slide.get("two_column"):
        blocks = "".join(
            f"<div><h2>{esc(h)}</h2><p>{esc(b)}</p></div>"
            for h, b in slide["two_column"]
        )
        parts.append(f'<div class="two-column">{blocks}</div>')
    if slide.get("risks"):
        items = "".join(f"<li>{esc(r)}</li>" for r in slide["risks"])
        parts.append(f'<ul class="risk-list">{items}</ul>')
    if slide.get("pre"):
        parts.append(f"<pre>{esc(slide['pre'])}</pre>")
    if slide.get("claim"):
        parts.append(f'<p class="claim">{esc(slide["claim"])}</p>')
    if slide.get("foot"):
        parts.append(f'<p class="foot">{esc(slide["foot"])}</p>')
    parts.append(f'<aside hidden>{esc(slide.get("notes", ""))}</aside>')
    return f'<section class="slide">{"".join(parts)}</section>'


def build_html() -> str:
    slides = "\n".join(render_slide(s) for s in SLIDES)
    return (
        '<!doctype html>\n<html lang="ko">\n<head>\n<meta charset="utf-8">\n'
        '<meta name="viewport" content="width=device-width, initial-scale=1">\n'
        "<title>FinLLM Lab — 발표자료</title>\n"
        f"<style>{CSS}</style>\n</head>\n<body>\n"
        f'<div class="deck">\n{slides}\n</div>\n</body>\n</html>\n'
    )


def build_pdf(html_path: Path, pdf_path: Path) -> bool:
    chrome = shutil.which("google-chrome") or shutil.which("chromium")
    if not chrome:
        print("WARNING: no chrome/chromium; skipping slide PDF", file=sys.stderr)
        return False
    completed = subprocess.run(
        [
            chrome, "--headless", "--disable-gpu", "--no-sandbox",
            "--no-pdf-header-footer",
            f"--print-to-pdf={pdf_path}", f"file://{html_path}",
        ],
        capture_output=True,
        text=True,
        timeout=180,
    )
    if completed.returncode != 0 or not pdf_path.exists():
        print(f"WARNING: chrome PDF failed: {completed.stderr[:300]}", file=sys.stderr)
        return False
    return True


def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Inches, Pt

    INK = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED = RGBColor(0x66, 0x66, 0x66)
    ACCENT = RGBColor(0x0B, 0x5C, 0xAD)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    SANS = "Malgun Gothic"
    MONO = "Consolas"

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    left, right = Inches(0.75), Inches(0.75)
    width = Inches(SLIDE_W_IN) - left - right

    def textbox(slide, top, height, text, size, *, bold=False, color=INK,
                font=SANS, align_center=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        para = frame.paragraphs[0]
        para.text = text
        run = para.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        if align_center:
            from pptx.enum.text import PP_ALIGN

            para.alignment = PP_ALIGN.CENTER
        return box

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE

        if spec.get("cover"):
            textbox(slide, Inches(2.3), Inches(1.0), spec["title"], 44,
                    bold=True, color=ACCENT)
            textbox(slide, Inches(3.3), Inches(0.7), spec["subtitle"], 24, bold=True)
            textbox(slide, Inches(4.3), Inches(1.6), "\n".join(spec["meta"]), 14,
                    color=MUTED)
            slide.notes_slide.notes_text_frame.text = spec.get("notes", "")
            continue

        cursor = Inches(0.55)
        textbox(slide, cursor, Inches(0.95), spec["title"], 27, bold=True)
        cursor += Inches(1.05)

        if spec.get("lead"):
            textbox(slide, cursor, Inches(0.6), spec["lead"], 15, bold=True)
            cursor += Inches(0.72)

        table_spec = spec.get("table")
        if table_spec:
            rows = len(table_spec["rows"]) + 1
            cols = len(table_spec["head"])
            height = Inches(0.34) * rows
            shape = slide.shapes.add_table(rows, cols, left, cursor, width, height)
            table = shape.table
            numeric = set(table_spec.get("numeric_cols", []))
            font_size = 10 if table_spec.get("compact") else 12
            for c, head in enumerate(table_spec["head"]):
                cell = table.cell(0, c)
                cell.text = str(head)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    para.runs[0].font.size = Pt(font_size)
                    para.runs[0].font.bold = True
                    para.runs[0].font.color.rgb = INK
                    para.runs[0].font.name = SANS
            for r, row in enumerate(table_spec["rows"], start=1):
                emphasis = table_spec.get("emphasis") == r - 1
                for c, value in enumerate(row):
                    cell = table.cell(r, c)
                    cell.text = str(value)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                    para = cell.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        run.font.size = Pt(font_size)
                        run.font.bold = emphasis
                        run.font.color.rgb = ACCENT if emphasis else INK
                        run.font.name = MONO if c in numeric else SANS
            cursor += height + Inches(0.25)

        if spec.get("sections"):
            count = len(spec["sections"])
            col_w = int((width - Emu(int(Inches(0.3)) * (count - 1))) / count)
            for i, (head, body) in enumerate(spec["sections"]):
                x = left + Emu(int(col_w + Inches(0.3)) * i)
                box = slide.shapes.add_textbox(x, cursor, Emu(col_w), Inches(1.8))
                frame = box.text_frame
                frame.word_wrap = True
                frame.paragraphs[0].text = head
                run = frame.paragraphs[0].runs[0]
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = ACCENT
                run.font.name = SANS
                para = frame.add_paragraph()
                para.text = body
                para.runs[0].font.size = Pt(12)
                para.runs[0].font.color.rgb = INK
                para.runs[0].font.name = SANS
            cursor += Inches(2.0)

        if spec.get("two_column"):
            col_w = int((width - Inches(0.5)) / 2)
            for i, (head, body) in enumerate(spec["two_column"]):
                x = left + Emu(int(col_w + Inches(0.5)) * i)
                box = slide.shapes.add_textbox(x, cursor, Emu(col_w), Inches(1.9))
                frame = box.text_frame
                frame.word_wrap = True
                frame.paragraphs[0].text = head
                run = frame.paragraphs[0].runs[0]
                run.font.size = Pt(15)
                run.font.bold = True
                run.font.color.rgb = ACCENT
                run.font.name = MONO
                para = frame.add_paragraph()
                para.text = body
                para.runs[0].font.size = Pt(12)
                para.runs[0].font.color.rgb = INK
                para.runs[0].font.name = SANS
            cursor += Inches(2.1)

        if spec.get("risks"):
            box = slide.shapes.add_textbox(left, cursor, width, Inches(4.0))
            frame = box.text_frame
            frame.word_wrap = True
            for i, risk in enumerate(spec["risks"]):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                para.text = f"· {risk}"
                para.runs[0].font.size = Pt(14)
                para.runs[0].font.color.rgb = INK
                para.runs[0].font.name = SANS
                para.space_after = Pt(10)
            cursor += Inches(3.0)

        if spec.get("pre"):
            box = slide.shapes.add_textbox(left, cursor, width, Inches(1.5))
            frame = box.text_frame
            frame.word_wrap = True
            for i, line in enumerate(spec["pre"].split("\n")):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                para.text = line
                if para.runs:
                    para.runs[0].font.size = Pt(10)
                    para.runs[0].font.color.rgb = INK
                    para.runs[0].font.name = MONO
            cursor += Inches(1.6)

        if spec.get("claim"):
            textbox(slide, cursor, Inches(0.75), spec["claim"], 15,
                    bold=True, color=ACCENT)

        if spec.get("foot"):
            textbox(slide, Inches(SLIDE_H_IN - 1.0), Inches(0.6), spec["foot"], 10,
                    color=MUTED)

        slide.notes_slide.notes_text_frame.text = spec.get("notes", "")

    prs.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--outdir", type=Path, default=OUT)
    args = parser.parse_args()
    args.outdir.mkdir(parents=True, exist_ok=True)

    html_path = args.outdir / "finllm-slides.html"
    html_path.write_text(build_html(), encoding="utf-8")
    print(f"{html_path}  ({html_path.stat().st_size / 1024:.0f} KB, {len(SLIDES)} slides)")

    pdf_path = args.outdir / "finllm-slides.pdf"
    if build_pdf(html_path, pdf_path):
        print(f"{pdf_path}  ({pdf_path.stat().st_size / 1024:.0f} KB)")

    pptx_path = args.outdir / "finllm-slides.pptx"
    build_pptx(pptx_path)
    print(f"{pptx_path}  ({pptx_path.stat().st_size / 1024:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
