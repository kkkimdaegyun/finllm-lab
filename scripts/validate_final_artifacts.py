#!/usr/bin/env python3
"""Deterministic evidence and artifact validation for the final handoff."""

from __future__ import annotations

import json
import math
import re
import shutil
import subprocess
import sys
from pathlib import Path

import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
REPORT = ROOT / "docs/final-report/final-report.md"
SLIDES = ROOT / "portfolio/slide-data.json"
PDF = ROOT / "artifacts/FinLLM-Lab-v0.2-Final-Technical-Report.pdf"
PPTX = ROOT / "artifacts/FinLLM-Lab-v0.2-Developer-Portfolio.pptx"
GATE = ROOT / "ops/evidence/final-rehearsal/gate-all.json"
REVIEW = ROOT / "docs/final-review/final-release-review.json"


def fail(message: str) -> None:
    raise AssertionError(message)


def approx(actual: float, expected: float, tolerance: float = 0.002) -> None:
    if not math.isclose(actual, expected, abs_tol=tolerance):
        fail(f"value mismatch: {actual} != {expected}")


def validate_evidence() -> None:
    result_paths = [
        ROOT / f"results/2026-08-08c-profile-a-qwen3-14b-awq-deploymentmatched-eager-r{i}.json"
        for i in (1, 2, 3)
    ]
    results = [json.loads(path.read_text(encoding="utf-8")) for path in result_paths]
    def mean(path):
        values=[]
        for payload in results:
            value=payload
            for key in path: value=value[key]
            values.append(float(value))
        return sum(values)/len(values)
    approx(mean(["metrics","quality_score"]),97.667,0.001)
    approx(mean(["metrics","p95_ttft_ms"]),129.995,0.001)
    approx(mean(["metrics","p95_user_ttft_ms"]),1273.402,0.001)
    approx(mean(["metrics","aggregate_output_tokens_per_s"]),315.331,0.001)
    approx(mean(["metrics","peak_vram_gib"]),21.961,0.001)
    for payload in results:
        if payload["evidence_type"] != "memory-budget-emulation": fail("evidence boundary changed")
        if payload["hardware"]["gpu_model"] != "NVIDIA RTX A6000": fail("GPU provenance changed")

    gate=json.loads(GATE.read_text(encoding="utf-8"))
    if gate["overall"]!="pass" or gate["counts"]!={"pass":11,"fail":0,"skipped":0}: fail("gate is not full PASS")
    smoke=next(stage for stage in gate["stages"] if stage["stage"]=="smoke-evaluation")
    approx(smoke["evidence"]["quality_score"],98.333,0.001)
    review=json.loads(REVIEW.read_text(encoding="utf-8"))
    if review["verdict"]!="PASS" or review["blockers"] or review["failed_gates"]: fail("final review inconsistent")


def validate_sources() -> None:
    report=REPORT.read_text(encoding="utf-8")
    slides=json.loads(SLIDES.read_text(encoding="utf-8"))
    text=report+json.dumps(slides,ensure_ascii=False)
    required=["memory-budget-emulation","NOT_EXECUTED","95dd24deba5669919e12b8535dbaf3128646ae5e","43.765","6.332","20/20","98.333"]
    for token in required:
        if token not in text: fail(f"required evidence token missing: {token}")
    forbidden=[r"RTX 4090[^\n]{0,40}315",r"native 24GB[^\n]{0,30}(VERIFIED|PASS)",r"FINAL RELEASE · FAIL"]
    for pattern in forbidden:
        if re.search(pattern,text,re.IGNORECASE): fail(f"forbidden claim: {pattern}")
    if len(slides["slides"])!=12: fail("portfolio must have 12 slides")


def validate_pdf() -> None:
    if not PDF.exists() or PDF.stat().st_size<100_000: fail("PDF missing or too small")
    document=pymupdf.open(PDF)
    if len(document)<10: fail("PDF unexpectedly short")
    extracted="\n".join(page.get_text() for page in document)
    for token in ("Executive Summary","FinLLM Lab v0.2","43.765","NOT_EXECUTED"):
        if token not in extracted: fail(f"PDF text missing: {token}")
    for page in document:
        if page.rect.width<=0 or page.rect.height<=0: fail("invalid PDF page")
    document.close()
    previews=sorted((ROOT/"portfolio/assets/previews/report").glob("page-*.png"))
    if len(previews)<10: fail("PDF render previews missing")


def validate_html_previews(slide_count: int) -> None:
    previews=sorted((ROOT/"portfolio/assets/previews/html").glob("slide-*.png"))
    if len(previews)!=slide_count: fail("HTML preview count mismatch")
    for path in previews:
        with Image.open(path) as image:
            if image.size!=(1280,720): fail(f"bad HTML render size: {path} {image.size}")
            extrema=image.convert("RGB").getextrema()
            if all(low==high for low,high in extrema): fail(f"blank HTML render: {path}")


def validate_pptx(data) -> None:
    if not PPTX.exists() or PPTX.stat().st_size<50_000: fail("PPTX missing or too small")
    prs=Presentation(PPTX)
    if len(prs.slides)!=len(data["slides"]): fail("PPTX slide count mismatch")
    ratio=prs.slide_width/prs.slide_height
    if not math.isclose(ratio,16/9,rel_tol=.001): fail("PPTX is not 16:9")
    for slide,item in zip(prs.slides,data["slides"]):
        texts="\n".join(shape.text for shape in slide.shapes if hasattr(shape,"text_frame"))
        if item["title"] not in texts: fail(f"PPTX title mismatch: {item['number']}")
        editable=sum(1 for shape in slide.shapes if shape.shape_type in {MSO_SHAPE_TYPE.TEXT_BOX,MSO_SHAPE_TYPE.AUTO_SHAPE,MSO_SHAPE_TYPE.TABLE})
        if editable<8: fail(f"too few editable objects on slide {item['number']}: {editable}")
        image_shapes=sum(shape.shape_type==MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        if image_shapes: fail(f"slide {item['number']} contains raster background/picture")


def validate_optional_pptx_render(slide_count: int) -> str:
    rendered=ROOT/"artifacts/rendered/FinLLM-Lab-v0.2-Developer-Portfolio.pdf"
    if rendered.exists():
        doc=pymupdf.open(rendered); count=len(doc); doc.close()
        if count!=slide_count: fail(f"rendered PPTX PDF page count {count} != {slide_count}")
        return f"VERIFIED ({count} pages through LibreOffice)"
    return "PENDING_VALIDATION (LibreOffice/PowerPoint render not present)"


def main() -> int:
    data=json.loads(SLIDES.read_text(encoding="utf-8"))
    validate_evidence(); validate_sources(); validate_pdf(); validate_html_previews(len(data["slides"])); validate_pptx(data)
    render_status=validate_optional_pptx_render(len(data["slides"]))
    print("evidence cross-check: PASS")
    print("PDF structure/text/raster render: PASS")
    print("HTML Chrome render (1280x720): PASS")
    print("PPTX editable OOXML/16:9/content parity: PASS")
    print("PPTX visual render:",render_status)
    return 0


if __name__ == "__main__":
    try: raise SystemExit(main())
    except AssertionError as error:
        print(f"FINAL ARTIFACT VALIDATION: FAIL — {error}",file=sys.stderr); raise SystemExit(1)
