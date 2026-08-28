#!/usr/bin/env python3
"""Build editable PPTX and Chrome previews from portfolio/slide-data.json."""

from __future__ import annotations

import json
import shutil
import subprocess
import threading
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PORTFOLIO = ROOT / "portfolio"
DATA_PATH = PORTFOLIO / "slide-data.json"
ARTIFACT = ROOT / "artifacts" / "FinLLM-Lab-v0.2-Developer-Portfolio.pptx"
HTML_ARTIFACT = ROOT / "artifacts" / "FinLLM-Lab-v0.2-Developer-Portfolio.html"
PREVIEW_DIR = PORTFOLIO / "assets" / "previews" / "html"
CONTACT = PORTFOLIO / "assets" / "previews" / "html-contact-sheet.png"
FONT = "Noto Sans KR"
SW, SH = 13.333, 7.5


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_shape(slide, x, y, w, h, fill, line=None, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w, h, size=16, color="#F5F8FC", bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    paragraph.space_after = Pt(0)
    return box


def add_rich_list(slide, values, x, y, w, h, size=13, color="#A4B5CA"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True
    frame.margin_left = Inches(.05); frame.margin_right = Inches(.02)
    for index, value in enumerate(values):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = str(value); p.level = 0; p.font.name = FONT; p.font.size = Pt(size)
        p.font.color.rgb = rgb(color); p.space_after = Pt(4); p.text = "• " + p.text
    return box


def base(slide, item, theme):
    background = slide.background.fill; background.solid(); background.fore_color.rgb = rgb(theme["background"])
    add_shape(slide, 0, 0, .09, SH, theme["cyan"], theme["cyan"], False)
    add_text(slide, item["eyebrow"], .62, .35, 10.7, .25, 8.5, theme["cyan"], True)
    add_text(slide, f'{item["number"]:02d}', 12.15, .31, .55, .3, 16, "#53708F", True, PP_ALIGN.RIGHT)
    add_text(slide, item["title"], .62, .68, 11.7, .58, 23, theme["text"], True)
    line = add_shape(slide, .62, 1.3, 12.05, .012, theme["line"], theme["line"], False)
    line.line.fill.background()
    add_shape(slide, .62, 7.05, 12.05, .012, theme["line"], theme["line"], False)
    add_text(slide, item["footer"], .62, 7.12, 6.5, .18, 7.2, theme["muted"])
    add_text(slide, "EVIDENCE · " + item["evidence"][0], 7.2, 7.12, 5.45, .18, 7.2, theme["muted"], False, PP_ALIGN.RIGHT)


def card(slide, value, x, y, w, h, theme, compact=False):
    add_shape(slide, x, y, w, h, theme["surface"], theme["line"])
    label = value.get("label") or value.get("severity") or ""
    add_text(slide, label, x+.15, y+.13, w-.3, .2, 7.5, theme["cyan"], True)
    add_text(slide, value.get("title", ""), x+.15, y+.39, w-.3, .3, 13 if not compact else 11.5, theme["text"], True)
    body = value.get("text") or value.get("status") or ""
    add_text(slide, body, x+.15, y+.75, w-.3, h-.88, 9.5 if not compact else 8.3, theme["muted"])
    if value.get("next"):
        add_text(slide, "Next · " + value["next"], x+.15, y+h-.3, w-.3, .17, 7.2, theme["amber"])


def render(slide, item, theme):
    base(slide, item, theme)
    c = item["content"]; layout = item["layout"]
    if layout == "cover":
        add_text(slide, c["headline"], .75, 1.75, 8.8, 1.45, 31, theme["text"], True)
        add_text(slide, c["subheadline"], .78, 3.35, 8.5, .38, 14, theme["muted"])
        x=.78
        for tag in c["tags"]:
            width=max(.72,len(tag)*.09+.3); add_shape(slide,x,3.95,width,.34,theme["surface"],theme["line"])
            add_text(slide,tag,x,4.02,width,.17,7.5,theme["text"],False,PP_ALIGN.CENTER); x+=width+.09
        add_shape(slide, 9.65, 4.75, 2.85, 1.03, "#12352F", theme["green"])
        add_text(slide,c["verdict"],9.83,4.97,2.5,.26,13,theme["green"],True)
        add_text(slide,c["verdict_note"],9.83,5.31,2.5,.3,8,theme["text"])
        return
    if layout in {"constraint-grid", "learning-grid", "risk-register"}:
        if c.get("headline"): add_text(slide,c["headline"],.72,1.52,11.9,.42,15,theme["text"])
        values=c.get("cards") or c.get("risks"); start_y=2.0 if c.get("headline") else 1.58
        columns=4 if layout=="learning-grid" else 3
        h=2.12 if layout=="learning-grid" else (1.75 if len(values)<=6 else 1.46)
        gap=.13; w=(11.9-(columns-1)*gap)/columns
        for i,v in enumerate(values): card(slide,v,.72+(i%columns)*(w+gap),start_y+(i//columns)*(h+gap),w,h,theme,compact=layout in {"learning-grid","risk-register"})
        if layout=="risk-register":
            add_shape(slide,.72,6.5,11.9,.38,"#12433E",theme["green"])
            add_text(slide,c["verdict"],.9,6.6,2.5,.18,10,theme["green"],True)
            add_text(slide,c["note"],3.1,6.59,9.2,.2,7.6,theme["text"])
        return
    if layout == "question":
        add_shape(slide,.85,1.65,11.6,1.15,theme["surface"],theme["cyan"])
        add_text(slide,c["quote"],1.15,1.92,11,.65,20,theme["text"],True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        w=3.73
        for i,v in enumerate(c["criteria"]):
            x=.85+(i%3)*(w+.2); y=3.08+(i//3)*1.02
            add_shape(slide,x,y,w,.83,theme["surface_alt"],theme["line"])
            add_text(slide,v["name"],x+.16,y+.12,w-.3,.2,10,theme["cyan"],True)
            add_text(slide,v["detail"],x+.16,y+.4,w-.3,.22,8.5,theme["muted"])
        add_text(slide,c["rule"],1,5.45,11.3,.35,11,theme["amber"],True,PP_ALIGN.CENTER)
        return
    if layout == "architecture":
        n=len(c["flow"]); w=2.05; gap=.33
        for i,v in enumerate(c["flow"]):
            x=.72+i*(w+gap); add_shape(slide,x,1.65,w,1.05,theme["surface"],theme["cyan"])
            add_text(slide,v["name"],x+.1,1.89,w-.2,.25,12,theme["text"],True,PP_ALIGN.CENTER)
            add_text(slide,v["detail"],x+.1,2.27,w-.2,.2,8,theme["muted"],False,PP_ALIGN.CENTER)
            if i<n-1:add_text(slide,"→",x+w+.04,2.02,.25,.3,18,theme["cyan"],True,PP_ALIGN.CENTER)
        for i,(title,values) in enumerate((("OBSERVE",c["ops"]),("CHANGE SAFETY",c["safety"]))):
            x=.72+i*6.03; add_shape(slide,x,3.13,5.85,1.53,theme["surface_alt"],theme["line"])
            add_text(slide,title,x+.18,3.32,2.2,.22,9,theme["amber"],True)
            add_rich_list(slide,values,x+.18,3.64,5.4,.68,9,theme["muted"])
        add_text(slide,c["status"],1,5.2,11.3,.34,11,theme["green"],True,PP_ALIGN.CENTER)
        return
    if layout == "data-table":
        rows=len(c["table"]["rows"])+1; cols=len(c["table"]["headers"])
        table=slide.shapes.add_table(rows,cols,Inches(.72),Inches(1.62),Inches(11.9),Inches(2.75)).table
        widths=[2.35,1.1,1.35,1.35,1.05,1.25,1.0]
        for i,width in enumerate(widths): table.columns[i].width=Inches(width)
        for col,value in enumerate(c["table"]["headers"]):
            cell=table.cell(0,col); cell.text=value; cell.fill.solid(); cell.fill.fore_color.rgb=rgb("#17395B")
        for r,row in enumerate(c["table"]["rows"],start=1):
            for col,value in enumerate(row):
                cell=table.cell(r,col); cell.text=value; cell.fill.solid(); cell.fill.fore_color.rgb=rgb("#14413F" if r-1==c["table"]["highlight_row"] else theme["surface"])
        for r in range(rows):
            for col in range(cols):
                cell=table.cell(r,col); cell.margin_left=cell.margin_right=Inches(.05); cell.margin_top=cell.margin_bottom=Inches(.04)
                for p in cell.text_frame.paragraphs:
                    p.font.name=FONT; p.font.size=Pt(8.5); p.font.bold=(r==0 or r-1==c["table"]["highlight_row"]); p.font.color.rgb=rgb(theme["cyan"] if r==0 else theme["text"]); p.alignment=PP_ALIGN.LEFT if col==0 else PP_ALIGN.RIGHT
        add_text(slide,c["note"],.78,4.58,11.7,.25,8.5,theme["muted"])
        add_shape(slide,.72,5.02,11.9,.85,"#15324B",theme["amber"])
        add_text(slide,c["decision"],.95,5.22,11.4,.4,10.5,theme["text"])
        return
    if layout == "autopsy":
        w=2.65
        for i,v in enumerate(c["steps"]):
            x=.72+i*3.03; add_shape(slide,x,1.62,w,2.03,theme["surface"],theme["line"])
            add_text(slide,v["label"],x+.15,1.78,w-.3,.18,8,theme["cyan"],True)
            add_text(slide,v["title"],x+.15,2.1,w-.3,.28,11.5,theme["text"],True)
            add_text(slide,v["metric"],x+.15,2.55,w-.3,.35,16,theme["amber"],True)
            add_text(slide,v["text"],x+.15,3.05,w-.3,.28,8.5,theme["muted"])
            if i<3:add_text(slide,"→",x+w+.12,2.4,.25,.3,18,theme["cyan"],True)
        add_shape(slide,.72,3.98,11.9,.85,"#15324B",theme["amber"])
        add_text(slide,c["corrected"],.95,4.18,11.4,.4,10,theme["text"])
        add_text(slide,c["principle"],1,5.22,11.3,.35,11,theme["amber"],True,PP_ALIGN.CENTER)
        return
    if layout == "evolution":
        for i,key in enumerate(("before","after")):
            v=c[key]; x=.72+i*6.25; add_shape(slide,x,1.6,5.65,2.1,theme["surface"],theme["line"])
            add_text(slide,v["label"],x+.22,1.8,2,.2,8,theme["cyan"],True)
            add_text(slide,v["title"],x+.22,2.15,5.1,.35,16,theme["text"],True)
            add_rich_list(slide,v["items"],x+.22,2.65,5,.72,9,theme["muted"])
        add_text(slide,"→",6.18,2.35,.5,.5,25,theme["cyan"],True,PP_ALIGN.CENTER)
        for i,v in enumerate(c["verification"]):
            x=.72+i*4.03; add_shape(slide,x,4.05,3.83,1.3,theme["surface_alt"],theme["line"])
            add_text(slide,v["status"],x+.15,4.2,.75,.18,7.5,theme["green"],True)
            add_text(slide,v["name"],x+.95,4.2,2.55,.18,9,theme["text"],True)
            add_text(slide,v["value"],x+.15,4.55,3.5,.28,14,theme["cyan"],True)
            add_text(slide,v["note"],x+.15,4.95,3.5,.2,7.5,theme["muted"])
        return
    if layout == "timeline":
        w=1.84
        for i,v in enumerate(c["timeline"]):
            x=.72+i*1.98; add_shape(slide,x,1.58,w,2.25,theme["surface"],theme["line"])
            add_shape(slide,x,1.58,w,.04,theme["cyan"],theme["cyan"],False)
            add_text(slide,v["time"],x+.12,1.78,w-.24,.18,7,theme["cyan"],True)
            add_text(slide,v["title"],x+.12,2.12,w-.24,.4,10,theme["text"],True)
            add_text(slide,v["value"],x+.12,2.65,w-.24,.35,14,theme["amber"],True)
            add_text(slide,v["detail"],x+.12,3.17,w-.24,.35,7.5,theme["muted"])
        add_shape(slide,.72,4.2,11.9,.68,"#15324B",theme["amber"])
        add_text(slide,c["lesson"],.95,4.37,11.4,.3,9.5,theme["text"])
        add_text(slide,c["caveat"],.8,5.15,11.7,.27,8.5,theme["muted"],False,PP_ALIGN.CENTER)
        return
    if layout == "review-loop":
        for i,v in enumerate(c["lanes"]):
            x=.72+i*6.05; add_shape(slide,x,1.58,5.85,1.48,theme["surface"],theme["line"])
            add_text(slide,v["owner"],x+.18,1.76,2.5,.2,8,theme["cyan"],True)
            add_text(slide,v["scope"],x+.18,2.08,5.4,.3,13,theme["text"],True)
            add_rich_list(slide,v["outputs"],x+.18,2.48,5.2,.33,8.5,theme["muted"])
        x=.72
        for i,v in enumerate(c["loop"]):
            width=1.65; add_shape(slide,x,3.38,width,.48,"#173B50",theme["line"])
            add_text(slide,v,x,3.53,width,.15,7.3,theme["text"],True,PP_ALIGN.CENTER); x+=width+.31
            if i<len(c["loop"])-1:add_text(slide,"→",x-.27,3.48,.24,.2,12,theme["cyan"],True)
        add_shape(slide,.72,4.18,7.65,1.25,theme["surface_alt"],theme["line"])
        add_rich_list(slide,c["findings"],.92,4.4,7.2,.75,8.5,theme["muted"])
        add_shape(slide,8.62,4.18,4,1.25,"#2D2A1A",theme["amber"])
        add_text(slide,c["judge"],8.88,4.48,3.5,.55,10,theme["amber"],True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        return
    if layout == "summary":
        add_shape(slide,.72,1.55,11.9,1.38,theme["surface"],theme["cyan"])
        add_text(slide,c["pitch"],1.02,1.82,11.3,.8,13,theme["text"],True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
        add_shape(slide,.72,3.22,7.7,1.7,theme["surface"],theme["line"])
        add_text(slide,"CORE STACK",.95,3.45,2,.2,8,theme["cyan"],True)
        x=.95; y=3.85
        for tag in c["stack"]:
            width=max(.7,len(tag)*.075+.25)
            if x+width>8.1: x=.95; y+=.42
            add_shape(slide,x,y,width,.31,theme["surface_alt"],theme["line"]); add_text(slide,tag,x,y+.08,width,.14,7,theme["text"],False,PP_ALIGN.CENTER); x+=width+.08
        add_shape(slide,8.65,3.22,3.97,1.7,theme["surface"],theme["line"])
        add_text(slide,c.get("side_label", "NEXT"),8.9,3.45,1.5,.2,8,theme["cyan"],True)
        add_rich_list(slide,c["next"],8.9,3.82,3.35,.72,9,theme["muted"])
        add_text(slide,c["closing"],.9,5.35,11.5,.35,10.5,theme["amber"],True,PP_ALIGN.CENTER)


def build_pptx(data):
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    blank=prs.slide_layouts[6]
    for item in data["slides"]:
        render(prs.slides.add_slide(blank),item,data["theme"])
    ARTIFACT.parent.mkdir(parents=True,exist_ok=True); prs.save(ARTIFACT)


def build_standalone_html(data):
    html=(PORTFOLIO / "index.html").read_text(encoding="utf-8")
    css=(PORTFOLIO / "styles.css").read_text(encoding="utf-8")
    payload=json.dumps(data,ensure_ascii=False,separators=(",", ":")).replace("</", "<\\/")
    html=html.replace('<link rel="stylesheet" href="styles.css">', f"<style>\n{css}\n</style>")
    html=html.replace("</head>", f'<script>window.__FINLLM_DATA__={payload};</script>\n</head>')
    HTML_ARTIFACT.parent.mkdir(parents=True,exist_ok=True)
    HTML_ARTIFACT.write_text(html,encoding="utf-8")


class QuietHandler(SimpleHTTPRequestHandler):
    def log_message(self, *_args): pass


def build_html_previews(slide_count: int):
    PREVIEW_DIR.mkdir(parents=True,exist_ok=True)
    for stale in PREVIEW_DIR.glob("slide-*.png"): stale.unlink()
    handler=lambda *args,**kwargs: QuietHandler(*args,directory=str(PORTFOLIO),**kwargs)
    server=ThreadingHTTPServer(("127.0.0.1",0),handler)
    thread=threading.Thread(target=server.serve_forever,daemon=True); thread.start()
    chrome=shutil.which("google-chrome") or shutil.which("chromium")
    if not chrome: raise SystemExit("Chrome/Chromium not found")
    try:
        for number in range(1,slide_count+1):
            target=PREVIEW_DIR/f"slide-{number:02d}.png"
            subprocess.run([chrome,"--headless=new","--no-sandbox","--disable-gpu","--hide-scrollbars","--force-device-scale-factor=1","--window-size=1280,720","--virtual-time-budget=2500",f"--screenshot={target}",f"http://127.0.0.1:{server.server_port}/index.html?slide={number}"],check=True,stdout=subprocess.DEVNULL,stderr=subprocess.DEVNULL)
    finally:
        server.shutdown(); server.server_close(); thread.join(timeout=2)
    thumbs=[]
    for path in sorted(PREVIEW_DIR.glob("slide-*.png")):
        image=Image.open(path).convert("RGB"); image.thumbnail((384,216)); thumbs.append(image.copy())
    cols=3; rows=(len(thumbs)+cols-1)//cols
    sheet=Image.new("RGB",(cols*410,rows*250),"#101b2a"); draw=ImageDraw.Draw(sheet)
    for i,image in enumerate(thumbs):
        x=(i%cols)*410+13; y=(i//cols)*250+27; sheet.paste(image,(x,y)); draw.text((x,y-19),f"Slide {i+1}",fill="white")
    CONTACT.parent.mkdir(parents=True,exist_ok=True); sheet.save(CONTACT)


def main() -> int:
    data=json.loads(DATA_PATH.read_text(encoding="utf-8"))
    if data["deck"]["aspect_ratio"]!="16:9": raise SystemExit("deck must be 16:9")
    build_standalone_html(data); build_pptx(data); build_html_previews(len(data["slides"]))
    print(f"HTML: {HTML_ARTIFACT} ({HTML_ARTIFACT.stat().st_size} bytes)")
    print(f"PPTX: {ARTIFACT} ({len(data['slides'])} editable slides, {ARTIFACT.stat().st_size} bytes)")
    print(f"HTML previews: {CONTACT}")
    return 0


if __name__ == "__main__": raise SystemExit(main())
