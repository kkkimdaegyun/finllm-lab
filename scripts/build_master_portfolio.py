#!/usr/bin/env python3
"""Build an editable HTML/PPTX portfolio covering all four developer projects."""

from __future__ import annotations

import html
import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ARTIFACTS = ROOT / "artifacts"
HTML_PATH = ARTIFACTS / "AI-Inference-Developer-Portfolio.html"
PPTX_PATH = ARTIFACTS / "AI-Inference-Developer-Portfolio.pptx"
DATA_PATH = ARTIFACTS / "AI-Inference-Developer-Portfolio.json"

FONT = "Noto Sans KR"
MONO = "DejaVu Sans Mono"
BG = "07111F"
PANEL = "0E1D30"
PANEL_ALT = "10263A"
TEXT = "F3F7FD"
MUTED = "9DB0C7"
CYAN = "38D6FF"
LIME = "9AF56D"
AMBER = "FFC857"
LINE = "203550"


SLIDES = [
    {
        "layout": "title",
        "kicker": "AI DEVELOPER · LLMOPS · AGENT ENGINEERING",
        "title": "Evidence-Driven\nAI Inference Portfolio",
        "lead": "김대균 · AI 추론을 만들고, 측정하고, 운영하며, 증거 기반 Agent로 자동화했습니다.",
        "chips": ["Production LLMOps", "Quantization Debugging", "Korean ASR", "FinCompliance Agent"],
    },
    {
        "layout": "cards",
        "kicker": "PORTFOLIO MAP",
        "title": "네 프로젝트가 답하는 질문",
        "lead": "모델 데모보다 요구사항 → 측정 → 선택 → 운영 → 자동화의 닫힌 고리를 만들었습니다.",
        "cards": [
            {"head": "01 FinLLM Lab", "value": "운영", "text": "금융 Private RAG를 단일 GPU에서 어떻게 배포·관측·회귀·롤백할까?", "accent": CYAN},
            {"head": "02 Quantization Autopsy", "value": "원인", "text": "양자화 모델의 성능 저하를 정밀도가 아닌 실행 경로까지 어떻게 분해할까?", "accent": LIME},
            {"head": "03 K-VoiceBench", "value": "평가", "text": "한국어 ASR을 정확도·엔티티·RTF·VRAM의 같은 계약으로 어떻게 비교할까?", "accent": AMBER},
            {"head": "04 FinCompliance Agent", "value": "행동", "text": "금융 준법 검토에서 모델이 도구를 선택하되 권한과 근거를 코드로 어떻게 강제할까?", "accent": CYAN},
        ],
    },
    {
        "layout": "flow",
        "kicker": "PROJECT 01 · FINLLM LAB V0.2",
        "title": "금융 RAG를 ‘운영 루프’로 완성",
        "lead": "ACL을 모델 앞에서 강제하고, 실제 장애와 회귀를 release gate까지 연결했습니다.",
        "flow": ["Request", "ACL Retrieval", "vLLM", "Metrics", "Regression Gate", "Rollback"],
        "metrics": [
            ["Final release", "PASS"], ["Deterministic tests", "153"],
            ["Actual gate", "11 / 11"], ["Result records", "27"],
        ],
        "callout": "Deployment → Observability → Alert → Incident → Rollback을 Compose rehearsal로 검증",
    },
    {
        "layout": "split",
        "kicker": "FINLLM · DEPLOYMENT DECISION",
        "title": "더 큰 14B가 더 적은 메모리로 합격",
        "lead": "4-bit 가중치 절감분을 KV cache로 돌려 품질과 동시성 여유를 확보했습니다.",
        "left": {
            "title": "채택 구성",
            "big": "Qwen3-14B-AWQ",
            "lines": ["gpu_memory_utilization 0.46", "--enforce-eager", "동시성 10 · 8,192 tokens"],
        },
        "right_metrics": [
            ["Quality", "97.667"], ["User P95 TTFT", "1,273 ms"],
            ["Throughput", "315.3 tok/s"], ["Peak VRAM", "21.96 GiB"],
        ],
        "callout": "A6000 memory-budget-emulation이며 실제 24GB 카드 성능은 NOT_EXECUTED",
    },
    {
        "layout": "bars",
        "kicker": "PROJECT 02 · QUANTIZATION AUTOPSY",
        "title": "‘4-bit라서 느리다’는 설명을 반증",
        "lead": "같은 모델·예산에서 실행 경로만 바꾸자 AWQ 처리량이 5.47배 회복됐습니다.",
        "bars": [
            ["8B BF16 · CUDA graph", 295.98, CYAN],
            ["8B BF16 · eager", 286.95, CYAN],
            ["14B AWQ · CUDA graph", 57.21, AMBER],
            ["14B AWQ · eager", 313.24, LIME],
        ],
        "metrics": [["Immutable runs", "18"], ["Configurations", "6"], ["Repeats", "3×"], ["AWQ recovery", "5.47×"]],
        "callout": "정확한 kernel root cause는 프로파일러로 측정하지 않아 NOT_MEASURED",
    },
    {
        "layout": "flow",
        "kicker": "QUANTIZATION · DEBUGGING METHOD",
        "title": "그럴듯한 원인보다 반증 가능한 실험",
        "lead": "비교 축을 고정하고 설명이 깨지는 조건을 찾았습니다.",
        "flow": ["관측\nAWQ 1/5", "가설\ndequant 비용", "통제\n같은 모델·예산", "반증\neager 회복", "제한\n원인 범위 명시"],
        "metrics": [
            ["Weight saved", "5.90 GiB"], ["KV cache gain", "5.58 GiB"],
            ["User TTFT ↓", "80.6%"], ["24GB headroom", "2.04 GiB"],
        ],
        "callout": "8B BF16과 14B AWQ는 크기가 달라 순수한 양자화 효과를 식별한 실험은 아님",
    },
    {
        "layout": "flow",
        "kicker": "PROJECT 03 · K-VOICEBENCH",
        "title": "한국어 ASR을 하나의 결과 계약으로 평가",
        "lead": "학습 없이 공개 Whisper 체크포인트의 추론 품질과 단일 GPU 운영성을 측정했습니다.",
        "flow": ["JSONL Manifest", "Audio Contract", "Whisper FP16", "WER·CER·Entity", "Latency·RTF·VRAM"],
        "metrics": [["GPU policy", "GPU 1 only"], ["Execution", "Sequential"], ["Training", "None"], ["CI backend", "Mock"]],
        "callout": "원본 음성·모델 가중치·로컬 절대 경로를 결과 JSON과 Git에서 분리",
    },
    {
        "layout": "table",
        "kicker": "K-VOICEBENCH · VERIFIED BASELINE",
        "title": "Whisper large-v3 · A6000 한 장 실측",
        "lead": "24개, 100.848초 금융 TTS 최소대립쌍 · clean 12 / telephony 12",
        "headers": ["지표", "전체", "Clean", "Telephony"],
        "rows": [
            ["ITN WER", "6.897%", "6.897%", "6.897%"],
            ["띄어쓰기 비민감 ITN CER", "0.000%", "0.000%", "0.000%"],
            ["금융 숫자 Entity 정확도", "100.0%", "100.0%", "100.0%"],
            ["P95 warm latency", "442.7ms", "—", "—"],
            ["Mean / P95 RTF", "0.0918 / 0.1017", "—", "—"],
            ["Peak allocated VRAM", "3,200.5MiB", "—", "—"],
        ],
        "callout": "TTS 진단값이며 사람 금융상담 음성 성능은 NOT_EVALUATED",
    },
    {
        "layout": "flow",
        "kicker": "PROJECT 04 · FINCOMPLIANCE AGENT",
        "title": "모델이 선택하고 애플리케이션이 통제",
        "lead": "LLM이 도구 결과를 관찰해 다음 행동을 고르지만 실행 권한은 policy engine이 강제합니다.",
        "flow": ["Request", "LLM Chooses", "Strict Policy", "Execute", "Observe·Replan", "HITL"],
        "metrics": [
            ["Strict tools", "9"], ["Policy eval", "12 / 12 PASS"],
            ["Automated tests", "12 / 12"], ["Unauthorized finalize", "0"],
        ],
        "callout": "실제 OpenAI model lane은 API key 미설정으로 NOT_EXECUTED · deterministic control lane 완료",
    },
    {
        "layout": "table",
        "kicker": "FINCOMPLIANCE · VERIFIED CONTROLS",
        "title": "Tool hallucination과 권한 경계를 회귀 테스트",
        "lead": "모델 선택을 결정적으로 재생하되 실제 tool registry·state·guardrail 경로를 그대로 통과시켰습니다.",
        "headers": ["항목", "구현 계약", "검증 결과", "Failure boundary"],
        "rows": [
            ["Tool call", "allowlist + strict JSON schema", "12 / 12 control PASS", "없는 도구·인자 차단"],
            ["Evidence", "returned citation allowlist", "grounded normal flow", "가짜 규정 ID 탐지"],
            ["Security", "untrusted output sanitizer", "injection neutralized", "도구 출력 명령 미실행"],
            ["Authority", "out-of-band single-use token", "unauthorized finalize 0", "모델이 승인 생성 불가"],
            ["Delivery", "FastAPI + Docker + CI + trace", "container smoke PASS", "trace_id replay"],
        ],
        "callout": "Control PASS는 금융 법률 정확도 100%가 아니라 expected block/recovery가 관찰됐다는 의미",
    },
    {
        "layout": "table",
        "kicker": "ENGINEERING CONTRACTS",
        "title": "재현성을 코드로 강제",
        "lead": "좋은 숫자보다 그 숫자가 어디서 왔는지 다시 확인할 수 있게 만들었습니다.",
        "headers": ["프로젝트", "입력 계약", "검증", "산출물"],
        "rows": [
            ["FinLLM Lab", "model·tokenizer revision, ACL corpus", "schema + 11-stage gate", "result JSON · ADR · incident"],
            ["Quantization Autopsy", "18 immutable runs", "3회 반복·provenance", "JSON · CSV · Markdown · HTML"],
            ["K-VoiceBench", "audio manifest·GPU visibility", "WER/CER/entity contract", "result JSON · HTML · CI"],
            ["FinCompliance Agent", "strict tool schema·synthetic evidence", "12 adversarial controls", "JSON trace · API · HTML · PPTX"],
        ],
        "callout": "미실행 항목은 PASS로 바꾸지 않고 NOT_EXECUTED / NOT_EVALUATED로 남김",
    },
    {
        "layout": "cards",
        "kicker": "DEVELOPER OWNERSHIP",
        "title": "제가 직접 닫은 문제",
        "lead": "구현 범위는 모델 호출이 아니라 데이터·서빙·측정·운영 자동화 전체입니다.",
        "cards": [
            {"head": "Inference", "value": "Runtime", "text": "vLLM·Transformers, single-GPU scheduling, memory budget, TTFT·RTF", "accent": CYAN},
            {"head": "Evaluation", "value": "Evidence", "text": "RAG quality·ACL·injection, WER·CER·entity, immutable provenance", "accent": LIME},
            {"head": "Operations", "value": "Safety", "text": "Docker Compose, metrics·alerts, regression gate, incident·rollback, CI", "accent": AMBER},
            {"head": "Agent", "value": "Tool Use", "text": "function calling loop, strict guardrails, prompt-injection defense, out-of-band HITL", "accent": CYAN},
        ],
    },
    {
        "layout": "table",
        "kicker": "VALIDATION STATUS",
        "title": "현재 완성 상태",
        "lead": "채용 담당자가 바로 열어보고, 개발자가 명령으로 다시 검증할 수 있습니다.",
        "headers": ["프로젝트", "검증 결과", "포트폴리오", "Git"],
        "rows": [
            ["FinLLM Lab v0.2", "153 tests + artifacts PASS", "HTML · PPTX · PDF", "origin/main"],
            ["Quantization Autopsy", "18 runs / 6 configs PASS", "Interactive HTML · report", "FinLLM companion"],
            ["K-VoiceBench", "6 tests + GPU1 result PASS", "Interactive HTML · evidence JSON", "local main"],
            ["FinCompliance Agent", "12 tests + 12/12 controls PASS", "HTML · PPTX · trace JSON", "public main"],
        ],
        "callout": "FinCompliance Agent는 public main · K-VoiceBench와 EvidenceOps는 원격 연결 대기",
    },
    {
        "layout": "title",
        "kicker": "ONE NARRATIVE",
        "title": "Build · Measure · Operate · Act",
        "lead": "금융 LLMOps, 양자화 원인 분석, 한국어 음성 평가를 거쳐 모델이 근거와 권한 안에서 행동하는 Agent를 구현했습니다.",
        "chips": ["AI Developer", "LLMOps Engineer", "Inference Engineer", "Agent Engineer"],
    },
]


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_font(run, size: float, color: str = TEXT, bold: bool = False, mono: bool = False) -> None:
    family = MONO if mono else FONT
    run.font.name = family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    props = run._r.get_or_add_rPr()
    props.set("lang", "ko-KR")
    east = props.find(qn("a:ea"))
    if east is None:
        east = OxmlElement("a:ea")
        props.append(east)
    east.set("typeface", family)


def add_text(slide, x, y, w, h, text, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, mono=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(.02)
    frame.margin_top = frame.margin_bottom = Inches(.02)
    frame.vertical_anchor = valign
    for index, line in enumerate(str(text).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.05
        for run in paragraph.runs:
            set_font(run, size, color, bold, mono)
    return box


def add_rect(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    return shape


def add_header(slide, data, number):
    add_text(slide, .72, .38, 8.8, .28, data.get("kicker", ""), 10, CYAN, True)
    add_text(slide, .72, .76, 11.9, .58, data["title"], 29, TEXT, True)
    if data.get("lead"):
        add_text(slide, .72, 1.38, 11.9, .58, data["lead"], 15, MUTED)
    add_text(slide, 11.8, .38, .8, .25, f"{number:02d}", 10, MUTED, False, PP_ALIGN.RIGHT, mono=True)


def add_footer(slide, number):
    add_text(slide, .72, 7.12, 5.2, .18, "김대균 · AI Inference Portfolio", 8.5, MUTED)
    add_text(slide, 12.0, 7.12, .6, .18, str(number), 8.5, MUTED, False, PP_ALIGN.RIGHT, mono=True)


def add_metric(slide, x, y, w, label, value, accent=LIME):
    add_rect(slide, x, y, w, 1.0)
    add_text(slide, x + .16, y + .14, w - .32, .22, label, 10, MUTED)
    add_text(slide, x + .16, y + .43, w - .32, .36, value, 20, accent, True, mono=True)


def add_callout(slide, text, y=6.32, color=AMBER):
    shape = add_rect(slide, .72, y, 11.9, .55, BG, color)
    shape.line.width = Pt(1.2)
    add_text(slide, .94, y + .14, 11.45, .23, text, 11.5, color, True)


def render_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for number, data in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid(); background.fore_color.rgb = rgb(BG)
        layout = data["layout"]

        if layout == "title":
            add_text(slide, .78, .82, 11.7, .28, data["kicker"], 11, CYAN, True)
            add_text(slide, .78, 1.55, 11.8, 2.2, data["title"], 43, TEXT, True, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, .82, 4.22, 10.9, .78, data["lead"], 18, MUTED)
            x = .82
            for chip in data["chips"]:
                width = max(1.8, len(chip) * .105)
                add_rect(slide, x, 5.38, width, .43, BG, LINE)
                add_text(slide, x + .11, 5.51, width - .22, .17, chip, 9.5, LIME, True, PP_ALIGN.CENTER)
                x += width + .16
            add_text(slide, .82, 6.72, 4, .2, "2026 · Developer Portfolio", 9, MUTED, mono=True)
            continue

        add_header(slide, data, number)
        if layout == "cards":
            cards = data["cards"]
            if len(cards) == 4:
                positions = [(.72, 2.08), (6.77, 2.08), (.72, 4.08), (6.77, 4.08)]
                card_w, card_h = 5.85, 1.72
            else:
                positions = [(.72 + index * 4.03, 2.2) for index in range(len(cards))]
                card_w, card_h = 3.83, 3.55
            for card, (x, y) in zip(cards, positions):
                add_rect(slide, x, y, card_w, card_h, PANEL_ALT)
                add_text(slide, x + .23, y + .22, card_w - .46, .22, card["head"], 9.5, card["accent"], True)
                add_text(slide, x + .23, y + .58, 1.5, .34, card["value"], 19 if len(cards) == 4 else 27, TEXT, True)
                text_x = x + 1.75 if len(cards) == 4 else x + .23
                text_y = y + .53 if len(cards) == 4 else y + 1.63
                text_w = card_w - 1.98 if len(cards) == 4 else card_w - .53
                add_text(slide, text_x, text_y, text_w, .8 if len(cards) == 4 else 1.18, card["text"], 11.5 if len(cards) == 4 else 15, MUTED)
                add_rect(slide, x + .23, y + card_h - .25, .86, .04, card["accent"], card["accent"], False)
        elif layout == "flow":
            count = len(data["flow"])
            gap = .13
            width = (11.9 - gap * (count - 1)) / count
            for index, label in enumerate(data["flow"]):
                x = .72 + index * (width + gap)
                add_rect(slide, x, 2.16, width, 1.18, PANEL_ALT)
                add_text(slide, x + .08, 2.46, width - .16, .58, label, 12.5, TEXT, True,
                         PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
                if index < count - 1:
                    add_text(slide, x + width - .02, 2.61, .17, .24, "→", 14, CYAN, True, PP_ALIGN.CENTER)
            metric_w = (11.9 - .39) / 4
            for index, (label, value) in enumerate(data["metrics"]):
                add_metric(slide, .72 + index * (metric_w + .13), 4.06, metric_w, label, value)
            add_callout(slide, data["callout"])
        elif layout == "split":
            add_rect(slide, .72, 2.08, 5.2, 3.88, PANEL_ALT)
            left = data["left"]
            add_text(slide, 1.0, 2.42, 4.65, .25, left["title"], 11, CYAN, True)
            add_text(slide, 1.0, 3.02, 4.65, .65, left["big"], 25, TEXT, True)
            add_text(slide, 1.0, 4.02, 4.65, 1.1, "\n".join(left["lines"]), 14, MUTED, mono=True)
            for index, (label, value) in enumerate(data["right_metrics"]):
                x = 6.18 + (index % 2) * 3.18
                y = 2.08 + (index // 2) * 1.52
                add_metric(slide, x, y, 3.0, label, value)
            add_callout(slide, data["callout"])
        elif layout == "bars":
            max_value = max(item[1] for item in data["bars"])
            for index, (label, value, color) in enumerate(data["bars"]):
                y = 2.08 + index * .67
                add_text(slide, .78, y + .08, 3.1, .24, label, 11, MUTED)
                add_rect(slide, 3.98, y, 5.1, .4, BG, LINE)
                add_rect(slide, 4.0, y + .02, 5.06 * value / max_value, .36, color, color)
                add_text(slide, 9.28, y + .07, 1.05, .22, f"{value:.2f}", 11, TEXT, True, PP_ALIGN.RIGHT, mono=True)
            for index, (label, value) in enumerate(data["metrics"]):
                add_metric(slide, 10.52, 2.08 + index * .92, 2.1, label, value, CYAN)
            add_callout(slide, data["callout"])
        elif layout == "table":
            headers, rows = data["headers"], data["rows"]
            table_shape = slide.shapes.add_table(
                len(rows) + 1, len(headers), Inches(.72), Inches(2.08), Inches(11.9), Inches(3.92)
            )
            table = table_shape.table
            table.first_row = True
            for col in range(len(headers)):
                table.columns[col].width = Inches(11.9 / len(headers))
            for col, value in enumerate(headers):
                table.cell(0, col).text = value
            for row_index, row in enumerate(rows, 1):
                for col, value in enumerate(row):
                    table.cell(row_index, col).text = value
            for row_index in range(len(rows) + 1):
                for col in range(len(headers)):
                    cell = table.cell(row_index, col)
                    cell.fill.solid(); cell.fill.fore_color.rgb = rgb(PANEL_ALT if row_index == 0 else PANEL)
                    cell.margin_left = cell.margin_right = Inches(.09)
                    cell.margin_top = cell.margin_bottom = Inches(.05)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT if col == 0 else PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            set_font(run, 10.5, TEXT if row_index == 0 else (TEXT if col == 0 else LIME), row_index == 0, col > 0)
            add_callout(slide, data["callout"])
        add_footer(slide, number)

    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def render_html() -> None:
    def card_markup(card):
        accent = f"#{card['accent']}"
        return (f'<article class="card" style="--accent:{accent}"><small>{html.escape(card["head"])}</small>'
                f'<strong>{html.escape(card["value"])}</strong><p>{html.escape(card["text"])}</p></article>')

    rendered = []
    for number, data in enumerate(SLIDES, 1):
        body = ""
        if data["layout"] == "title":
            chips = "".join(f"<span>{html.escape(item)}</span>" for item in data["chips"])
            body = f'<div class="title-body"><h1>{html.escape(data["title"]).replace(chr(10), "<br>")}</h1><p>{html.escape(data["lead"])}</p><div class="chips">{chips}</div></div>'
        else:
            body += f'<header><div class="kicker">{html.escape(data["kicker"])}</div><h2>{html.escape(data["title"])}</h2><p>{html.escape(data.get("lead", ""))}</p></header>'
            layout = data["layout"]
            if layout == "cards":
                body += f'<div class="cards count-{len(data["cards"])}">' + "".join(card_markup(card) for card in data["cards"]) + '</div>'
            elif layout == "flow":
                body += f'<div class="flow" style="grid-template-columns:repeat({len(data["flow"])},1fr)">' + "".join(f'<div>{html.escape(item).replace(chr(10), "<br>")}</div>' for item in data["flow"]) + '</div>'
                body += '<div class="metrics">' + "".join(f'<div><small>{html.escape(label)}</small><strong>{html.escape(value)}</strong></div>' for label, value in data["metrics"]) + '</div>'
            elif layout == "split":
                left = data["left"]
                body += f'<div class="split"><div class="focus"><small>{html.escape(left["title"])}</small><strong>{html.escape(left["big"])}</strong><p>{"<br>".join(html.escape(item) for item in left["lines"])}</p></div><div class="metrics">' + "".join(f'<div><small>{html.escape(label)}</small><strong>{html.escape(value)}</strong></div>' for label, value in data["right_metrics"]) + '</div></div>'
            elif layout == "bars":
                maximum = max(item[1] for item in data["bars"])
                body += '<div class="bar-layout"><div class="bars">' + "".join(f'<div class="bar-row"><span>{html.escape(label)}</span><i><b style="width:{value/maximum*100:.1f}%;background:#{color}"></b></i><em>{value:.2f}</em></div>' for label, value, color in data["bars"]) + '</div><div class="side-metrics">' + "".join(f'<div><small>{html.escape(label)}</small><strong>{html.escape(value)}</strong></div>' for label, value in data["metrics"]) + '</div></div>'
            elif layout == "table":
                head = "".join(f"<th>{html.escape(item)}</th>" for item in data["headers"])
                rows = "".join("<tr>" + "".join(f"<td>{html.escape(str(item))}</td>" for item in row) + "</tr>" for row in data["rows"])
                body += f"<table><thead><tr>{head}</tr></thead><tbody>{rows}</tbody></table>"
            if data.get("callout"):
                body += f'<div class="callout">{html.escape(data["callout"])}</div>'
        rendered.append(f'<section class="slide"><div class="slide-no">{number:02d}</div>{body}</section>')

    document = f"""<!doctype html><html lang="ko"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>AI Inference Developer Portfolio</title><style>
:root{{--bg:#07111f;--panel:#0e1d30;--panel2:#10263a;--text:#f3f7fd;--muted:#9db0c7;--cyan:#38d6ff;--lime:#9af56d;--amber:#ffc857;--line:#203550}}*{{box-sizing:border-box}}body{{margin:0;background:#040a12;color:var(--text);font-family:"Noto Sans KR","Noto Sans",sans-serif}}.deck{{display:grid;gap:28px;justify-content:center;padding:28px}}.slide{{width:1280px;height:720px;padding:42px 68px;background:radial-gradient(circle at 90% 0,#12345a 0,transparent 36%),var(--bg);position:relative;overflow:hidden}}.slide-no{{position:absolute;right:68px;top:40px;color:var(--muted);font:12px monospace}}.kicker{{color:var(--cyan);font-weight:800;letter-spacing:.14em;font-size:12px}}h2{{font-size:40px;letter-spacing:-.04em;margin:12px 0 10px}}header>p{{color:var(--muted);font-size:18px;margin:0 0 26px}}.title-body{{height:100%;display:flex;flex-direction:column;justify-content:center}}.title-body h1{{font-size:72px;line-height:.98;letter-spacing:-.055em;margin:0 0 28px}}.title-body p{{color:var(--muted);font-size:22px;max-width:980px;line-height:1.6}}.chips{{display:flex;gap:12px;margin-top:35px}}.chips span{{border:1px solid var(--line);border-radius:999px;padding:10px 15px;color:var(--lime)}}.cards{{display:grid;grid-template-columns:repeat(3,1fr);gap:18px}}.cards.count-4{{grid-template-columns:repeat(2,1fr)}}.card,.metrics>div,.side-metrics>div,.focus{{background:linear-gradient(145deg,var(--panel2),var(--panel));border:1px solid var(--line);border-radius:16px;padding:22px}}.card small,.metrics small,.side-metrics small,.focus small{{color:var(--accent,var(--muted));font-weight:700}}.card strong,.focus strong{{display:block;font-size:30px;margin:25px 0}}.card p,.focus p{{color:var(--muted);font-size:17px;line-height:1.55}}.flow{{display:grid;gap:10px}}.flow div{{background:var(--panel2);border:1px solid var(--line);border-radius:13px;padding:22px 10px;text-align:center;font-weight:700;min-width:0}}.metrics{{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin-top:22px}}.metrics strong,.side-metrics strong{{display:block;color:var(--lime);font:700 24px monospace;margin-top:10px}}.callout{{position:absolute;left:68px;right:68px;bottom:45px;border:1px solid var(--amber);border-radius:12px;padding:13px 18px;color:var(--amber);font-weight:700}}.split{{display:grid;grid-template-columns:1fr 1.2fr;gap:18px}}.split .metrics{{margin:0;grid-template-columns:1fr 1fr}}.bar-layout{{display:grid;grid-template-columns:4fr 1fr;gap:28px}}.bar-row{{display:grid;grid-template-columns:240px 1fr 70px;gap:12px;align-items:center;margin:15px 0;color:var(--muted)}}.bar-row i{{height:17px;background:#050d18;border-radius:999px;overflow:hidden}}.bar-row b{{display:block;height:100%;border-radius:999px}}.bar-row em{{font:normal 13px monospace;color:var(--text)}}.side-metrics{{display:grid;gap:9px}}table{{width:100%;border-collapse:collapse;background:rgba(14,29,48,.8);font-size:15px}}th,td{{padding:13px 14px;border-bottom:1px solid var(--line);text-align:center}}th:first-child,td:first-child{{text-align:left}}th{{color:var(--cyan)}}td:not(:first-child){{font-family:monospace;color:var(--lime)}}@media print{{.deck{{display:block;padding:0}}.slide{{break-after:page}}}}@media(max-width:900px){{.deck{{display:block;padding:0}}.slide{{width:100vw;height:auto;min-height:100vh}}.cards,.metrics,.flow{{grid-template-columns:1fr 1fr!important}}}}
</style></head><body><main class="deck">{"".join(rendered)}</main></body></html>"""
    HTML_PATH.write_text(document, encoding="utf-8")


def validate() -> None:
    presentation = Presentation(PPTX_PATH)
    if len(presentation.slides) != len(SLIDES):
        raise ValueError("slide count mismatch")
    for number, slide in enumerate(presentation.slides, 1):
        if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
            raise ValueError(f"slide {number}: raster picture found")
        if sum(hasattr(shape, "text_frame") for shape in slide.shapes) < 5:
            raise ValueError(f"slide {number}: too few editable shapes")
    with zipfile.ZipFile(PPTX_PATH) as archive:
        xml = "\n".join(
            archive.read(name).decode("utf-8", errors="ignore")
            for name in archive.namelist() if name.endswith(".xml")
        )
    if FONT not in xml:
        raise ValueError("Noto Sans KR font metadata missing")
    page = HTML_PATH.read_text(encoding="utf-8")
    for token in ("FinLLM", "Quantization Autopsy", "K-VoiceBench", "FinCompliance Agent", "NOT_EVALUATED"):
        if token not in page:
            raise ValueError(f"HTML token missing: {token}")


def main() -> None:
    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    DATA_PATH.write_text(json.dumps(SLIDES, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    render_html()
    render_pptx()
    validate()
    print(f"HTML -> {HTML_PATH}")
    print(f"PPTX -> {PPTX_PATH}")
    print(f"slides: {len(SLIDES)}, editable/font validation: PASS")


if __name__ == "__main__":
    main()
